import os
import numpy as np
from pptx import Presentation
from sentence_transformers import SentenceTransformer, util

def extract_slides_text(ppt_path):
    """Extract text from each slide of a PowerPoint presentation."""
    prs = Presentation(ppt_path)
    slides_text = []
    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        slides_text.append(" ".join(slide_text))  # Combine text in a slide
    return slides_text  

def extract_text_from_ppt(ppt_path):
    """Extract text from a PowerPoint presentation."""
    prs = Presentation(ppt_path)
    text_content = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_content.append(shape.text)
    return " ".join(text_content)  
    
# Load sentence transformer model
model = SentenceTransformer("all-MiniLM-L6-v2")

# Query PPT
query_ppt_path = "ppts/Future_of_AI_2.pptx"
query_slides = extract_slides_text(query_ppt_path)

# Top 4 similar PPTs
ppt_folder = "ppts"
ppt_files = [os.path.join(ppt_folder, f) for f in os.listdir(ppt_folder) if f.endswith(".pptx")][:10]

# Extract text from all PPTs and embed them
all_docs = [extract_text_from_ppt(ppt) for ppt in ppt_files]
embeddings = model.encode(all_docs, convert_to_tensor=True) 

# Encode the query document
query_embedding = model.encode(extract_text_from_ppt(query_ppt_path), convert_to_tensor=True)

# Compute similarity scores and Get top-k similar PPTs
similarity_scores = util.pytorch_cos_sim(query_embedding, embeddings)[0].cpu().numpy()
top_k = 4
top_indices = np.argsort(similarity_scores)[-top_k:][::-1]  # Get top 4 in descending order
top_ppts = [ppt_files[idx] for idx in top_indices]

print("Top 4 similar PPTs:")
for idx in top_indices:
    print(f"{ppt_files[idx]} - Similarity Score: {similarity_scores[idx]:.4f}")



top_ppts_slides = {ppt: extract_slides_text(ppt) for ppt in top_ppts}
query_slide_embeddings = model.encode(query_slides, convert_to_tensor=True)

# Compare each query slide with each slide from top 4 PPTs
for ppt, slides in top_ppts_slides.items():
    print(f"\nAnalyzing similarity with: {ppt}")
    ppt_slide_embeddings = model.encode(slides, convert_to_tensor=True)

    similarity_matrix = util.pytorch_cos_sim(query_slide_embeddings, ppt_slide_embeddings).cpu().numpy()

    for i, query_slide in enumerate(query_slides):
        top_slide_indices = np.argsort(similarity_matrix[i])[-3:][::-1]  # Top 3 similar slides
        print(f"\nQuery Slide {i+1}:")
        print(query_slide[:200] + "..." if len(query_slide) > 200 else query_slide)
        for j in top_slide_indices:
            print(f"  -> Similar Slide {j+1} in {ppt} - Score: {similarity_matrix[i][j]:.4f}")
            print(f"     {slides[j][:200] + '...' if len(slides[j]) > 200 else slides[j]}")
