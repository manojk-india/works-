from pptx import Presentation
from sentence_transformers import SentenceTransformer, util

# Function to extract text from PPT
def extract_text_from_ppt(ppt_path):
    prs = Presentation(ppt_path)
    text_content = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_content.append(shape.text)
    return " ".join(text_content)

# Load the pre-trained embedding model
model = SentenceTransformer("all-MiniLM-L6-v2")  # Efficient and accurate model

# Extract text from both PPTs
ppt1_path = "Future_of_AI_1.pptx"  # Update with correct path
ppt2_path = "Future_of_AI_2.pptx"

text1 = extract_text_from_ppt(ppt1_path)
text2 = extract_text_from_ppt(ppt2_path)

# Compute embeddings
embedding1 = model.encode(text1, convert_to_tensor=True)
embedding2 = model.encode(text2, convert_to_tensor=True)

# Compute cosine similarity
similarity_score = util.pytorch_cos_sim(embedding1, embedding2).item()

print(f"Similarity between PPTs: {similarity_score:.4f}")
